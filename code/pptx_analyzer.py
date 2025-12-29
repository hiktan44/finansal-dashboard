#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
PPTX Analiz Aracı - TÜRKİYE EKONOMİSİ ANALİZİ
PowerPoint dosyasındaki tüm içeriği analiz eder
"""

from pptx import Presentation
import json
import os
import re
from pathlib import Path

def analyze_pptx(file_path):
    """PPTX dosyasını analiz eder"""
    
    print(f"🔍 Analiz başlatılıyor: {file_path}")
    
    # PowerPoint dosyasını aç
    prs = Presentation(file_path)
    
    analysis_result = {
        "total_slides": len(prs.slides),
        "slides": [],
        "static_content": [],
        "charts_tables": [],
        "links": [],
        "website_bands": []
    }
    
    print(f"📊 Toplam slide sayısı: {len(prs.slides)}")
    
    # Her slide'ı analiz et
    for slide_num, slide in enumerate(prs.slides, 1):
        print(f"\n📄 Slide {slide_num} analiz ediliyor...")
        
        slide_data = {
            "slide_number": slide_num,
            "title": "",
            "content": [],
            "links": [],
            "charts": [],
            "tables": []
        }
        
        # Slide'daki tüm shape'leri kontrol et
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                text_content = shape.text.strip()
                
                # Title kontrolü
                if shape.shape_type == 1:  # Title placeholder
                    slide_data["title"] = text_content
                    print(f"  📋 Başlık: {text_content[:50]}...")
                
                # Genel içerik
                else:
                    slide_data["content"].append(text_content)
                    print(f"  📝 İçerik: {text_content[:50]}...")
                
                # Link tespiti (URL pattern)
                url_pattern = r'https?://[^\s<>"{}|\\^`\[\]]+'
                links_found = re.findall(url_pattern, text_content)
                if links_found:
                    slide_data["links"].extend(links_found)
                    analysis_result["links"].extend([(slide_num, link) for link in links_found])
                    print(f"  🔗 Link bulundu: {links_found}")
                
                # Email pattern
                email_pattern = r'\b[A-Za-z0-9._%+-]+@[A-Za-z0-9.-]+\.[A-Z|a-z]{2,}\b'
                emails_found = re.findall(email_pattern, text_content)
                if emails_found:
                    slide_data["links"].extend([f"mailto:{email}" for email in emails_found])
                    print(f"  📧 Email bulundu: {emails_found}")
        
        # Grafik/tablo tespiti
        for shape in slide.shapes:
            if shape.shape_type == 3:  # Chart
                slide_data["charts"].append({
                    "type": "chart",
                    "name": shape.name if hasattr(shape, 'name') else "Unknown Chart"
                })
                analysis_result["charts_tables"].append({
                    "slide": slide_num,
                    "type": "chart",
                    "name": shape.name if hasattr(shape, 'name') else "Unknown Chart"
                })
                print(f"  📊 Grafik tespit edildi: {shape.name}")
            
            elif shape.shape_type == 19:  # Table
                slide_data["tables"].append({
                    "type": "table", 
                    "name": shape.name if hasattr(shape, 'name') else "Unknown Table"
                })
                analysis_result["charts_tables"].append({
                    "slide": slide_num,
                    "type": "table",
                    "name": shape.name if hasattr(shape, 'name') else "Unknown Table"
                })
                print(f"  📋 Tablo tespit edildi: {shape.name}")
        
        analysis_result["slides"].append(slide_data)
    
    # Statik içerik analizi
    print(f"\n📈 Statik içerik analizi...")
    all_content = " ".join([slide["title"] for slide in analysis_result["slides"]] + 
                          [content for slide in analysis_result["slides"] for content in slide["content"]])
    
    # Ekonomi ile ilgili anahtar kelimeler
    economy_keywords = [
        "ekonomi", "enflasyon", "faiz", "GSYH", "işsizlik", "bütçe", 
        "cari açık", "dış ticaret", " TCMB", "TÜİK", "Merkez Bankası",
        "alternatif senaryo", "öngörü", "projeksiyon", "tahmin"
    ]
    
    for keyword in economy_keywords:
        if keyword.lower() in all_content.lower():
            analysis_result["static_content"].append(keyword)
    
    print(f"✅ Anahtar ekonomi kelimeleri: {analysis_result['static_content']}")
    
    return analysis_result

def generate_summary_report(analysis_result):
    """Analiz sonuçlarından özet rapor oluşturur"""
    
    print("\n" + "="*60)
    print("📊 PPTX ANALİZ RAPORU")
    print("="*60)
    
    print(f"📄 Toplam Slide: {analysis_result['total_slides']}")
    print(f"🔗 Toplam Link: {len(analysis_result['links'])}")
    print(f"📊 Grafik/Tablo: {len(analysis_result['charts_tables'])}")
    print(f"📈 Ekonomi Kelimeleri: {len(analysis_result['static_content'])}")
    
    print(f"\n🔗 BULUNAN LİNKLER:")
    for slide_num, link in analysis_result['links']:
        print(f"  Slide {slide_num}: {link}")
    
    print(f"\n📊 GRAFİK/TABLOLAR:")
    for item in analysis_result['charts_tables']:
        print(f"  Slide {item['slide']}: {item['type']} - {item['name']}")
    
    print(f"\n📈 ANAHTAR KELİMELER:")
    for keyword in analysis_result['static_content']:
        print(f"  • {keyword}")

def main():
    # PPTX dosya yolu
    pptx_file = "user_input_files/TÜRKİYE  EKONOMİSİNİN DEĞERLENDİRİLMESİ VE ÖNGÖRÜLER 5 Temmuz 2025dip notlu .pptx"
    
    try:
        # Analiz çalıştır
        result = analyze_pptx(pptx_file)
        
        # Özet rapor oluştur
        generate_summary_report(result)
        
        # Sonuçları JSON dosyasına kaydet
        output_file = "data/pptx_analysis_result.json"
        os.makedirs("data", exist_ok=True)
        
        with open(output_file, 'w', encoding='utf-8') as f:
            json.dump(result, f, ensure_ascii=False, indent=2)
        
        print(f"\n💾 Detaylı analiz kaydedildi: {output_file}")
        print("✅ PPTX analizi tamamlandı!")
        
        return result
        
    except Exception as e:
        print(f"❌ Hata oluştu: {str(e)}")
        return None

if __name__ == "__main__":
    main()